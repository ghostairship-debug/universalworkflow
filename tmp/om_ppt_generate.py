from __future__ import annotations

import json
from pathlib import Path
from shutil import copyfile

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SOURCE_DIR = Path(r"C:\Users\74755\Desktop\OM PPT")
ROOT = Path(r"D:\Universal Agentic workflow")
VERTEX_SPEC_PATH = ROOT / "state" / "om_ppt_slide_spec.json"
OUTPUT_DIR = SOURCE_DIR
ASSET_DIR = SOURCE_DIR / "vertex_assets"
OUTPUT_PATH = OUTPUT_DIR / "Cainiao_Operational_Management_Refined_Deck.pptx"

PRIMARY = RGBColor(0x0E, 0x1B, 0x2B)
SECONDARY = RGBColor(0x1B, 0x26, 0x38)
ACCENT = RGBColor(0xF6, 0x73, 0x1C)
ACCENT_SOFT = RGBColor(0xFF, 0xE2, 0xD1)
TEXT = RGBColor(0x16, 0x1F, 0x2C)
MUTED = RGBColor(0x5D, 0x6B, 0x7A)
BG = RGBColor(0xF4, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_LINE = RGBColor(0xD9, 0xE1, 0xE8)
GREEN = RGBColor(0x1E, 0x8E, 0x6D)
YELLOW = RGBColor(0xF1, 0xB5, 0x4C)


def copy_cover_image() -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    source = SOURCE_DIR / "vertex_probe.png"
    target = ASSET_DIR / "cover_hero.png"
    if source.exists():
        copyfile(source, target)
    return target


def load_vertex_spec() -> list[dict]:
    if not VERTEX_SPEC_PATH.exists():
        return []
    try:
        data = json.loads(VERTEX_SPEC_PATH.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return []
    if isinstance(data, list):
        return data
    return data.get("slides", [])


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, font_size=18, color=TEXT, bold=False,
                font_name="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
                margin_left=0.06, margin_right=0.06, margin_top=0.04, margin_bottom=0.02):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    box.text_frame.margin_left = Inches(margin_left)
    box.text_frame.margin_right = Inches(margin_right)
    box.text_frame.margin_top = Inches(margin_top)
    box.text_frame.margin_bottom = Inches(margin_bottom)
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = align
    if not p.runs:
        return box
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def add_bullets(slide, left, top, width, height, bullets, *, font_size=18, color=TEXT, bullet_color=ACCENT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.bullet = True
    return box


def add_title(slide, title, subtitle=None, *, dark=False):
    title_color = WHITE if dark else PRIMARY
    subtitle_color = RGBColor(0xDF, 0xE5, 0xEC) if dark else MUTED
    add_textbox(slide, Inches(0.65), Inches(0.36), Inches(8.8), Inches(0.68), title,
                font_size=28, color=title_color, bold=True, valign=MSO_ANCHOR.BOTTOM)
    if subtitle:
        add_textbox(slide, Inches(0.65), Inches(0.98), Inches(8.9), Inches(0.52), subtitle,
                    font_size=12.5, color=subtitle_color)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.65), Inches(1.46), Inches(0.9), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.color.rgb = ACCENT
    return line


def add_footer(slide, text="", number=None, *, dark=False):
    if text:
        add_textbox(slide, Inches(0.65), Inches(7.08), Inches(8.6), Inches(0.24), text,
                    font_size=8.5, color=RGBColor(0xBE, 0xC8, 0xD2) if dark else MUTED)
    if number is not None:
        add_textbox(slide, Inches(12.2), Inches(7.02), Inches(0.45), Inches(0.26), str(number),
                    font_size=10, color=RGBColor(0xBE, 0xC8, 0xD2) if dark else MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, text, *, fill=WHITE, title_color=PRIMARY, text_color=MUTED,
             border=LIGHT_LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(1.1)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.08), width - Inches(0.2), Inches(0.28), title,
                font_size=11.5, color=title_color, bold=True)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.34), width - Inches(0.2), height - Inches(0.42), text,
                font_size=15.5, color=text_color, bold=False, valign=MSO_ANCHOR.MIDDLE)
    return shape


def add_metric_band(slide, items, *, top=5.68):
    left = 0.65
    gap = 0.14
    width = (12.03 - (gap * (len(items) - 1))) / len(items)
    for idx, (label, value, tone) in enumerate(items):
        x = left + idx * (width + gap)
        fill = WHITE if tone == "default" else ACCENT_SOFT
        border = LIGHT_LINE if tone == "default" else ACCENT
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(width), Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = border
        shape.line.width = Pt(1.0)
        add_textbox(slide, Inches(x + 0.12), Inches(top + 0.1), Inches(width - 0.24), Inches(0.22), label,
                    font_size=10, color=MUTED, bold=True)
        add_textbox(slide, Inches(x + 0.12), Inches(top + 0.37), Inches(width - 0.24), Inches(0.4), value,
                    font_size=18, color=PRIMARY, bold=True, valign=MSO_ANCHOR.MIDDLE)


def add_chip(slide, left, top, width, text, *, fill=ACCENT_SOFT, text_color=PRIMARY, line=ACCENT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    add_textbox(slide, left + Inches(0.06), top + Inches(0.03), width - Inches(0.12), Inches(0.18), text,
                font_size=10.2, color=text_color, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_section_label(slide, left, top, text):
    add_textbox(slide, left, top, Inches(2.2), Inches(0.24), text,
                font_size=10.5, color=ACCENT, bold=True)


def slide1(prs: Presentation, hero: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if hero.exists():
        slide.shapes.add_picture(str(hero), 0, 0, width=prs.slide_width, height=prs.slide_height)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = PRIMARY
    overlay.fill.transparency = 0.38
    overlay.line.color.rgb = PRIMARY
    add_textbox(slide, Inches(0.72), Inches(4.75), Inches(8.8), Inches(1.28),
                "How Cainiao Builds a Smart Logistics Ecosystem",
                font_size=27, color=WHITE, bold=True, valign=MSO_ANCHOR.BOTTOM)
    add_textbox(slide, Inches(0.72), Inches(6.02), Inches(7.3), Inches(0.5),
                "An Operations Management case on platform coordination, standardization, and smarter logistics decisions.",
                font_size=13.2, color=RGBColor(0xE4, 0xEA, 0xF0))
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(4.45), Inches(1.2), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.color.rgb = ACCENT
    add_textbox(slide, Inches(0.72), Inches(6.76), Inches(5.8), Inches(0.22),
                "Operational Management | Cainiao case study | 14-slide refined deck",
                font_size=9.5, color=RGBColor(0xD5, 0xDE, 0xE7))
    add_footer(slide, "Vertex-assisted hero visual; structure rebuilt from the brief, not from the draft layout.", 1, dark=True)


def slide2(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "Why this case matters for operations management",
              "The real question is how a platform can improve speed, cost, visibility, flexibility, and scalability in a fragmented logistics system.")
    add_section_label(slide, Inches(0.7), Inches(1.82), "CORE QUESTION")
    add_textbox(slide, Inches(0.7), Inches(2.08), Inches(5.55), Inches(1.18),
                "How can a platform model improve logistics performance without owning every logistics asset? Cainiao's answer is to standardize information, orchestrate partners, and optimize decisions across the network.",
                font_size=18, color=TEXT)
    callout = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.65), Inches(1.98), Inches(5.05), Inches(1.38))
    callout.fill.solid()
    callout.fill.fore_color.rgb = WHITE
    callout.line.color.rgb = ACCENT
    callout.line.width = Pt(1.5)
    add_textbox(slide, Inches(6.9), Inches(2.22), Inches(4.55), Inches(0.22), "Thesis", font_size=11, color=ACCENT, bold=True)
    add_textbox(slide, Inches(6.9), Inches(2.48), Inches(4.5), Inches(0.62),
                "Cainiao improves logistics by coordinating information first, so physical assets can move more predictably and at lower coordination cost.",
                font_size=16.5, color=PRIMARY)
    strip_y = 4.82
    pillars = [
        ("Speed", "Compress cycle time"),
        ("Cost", "Reduce wasteful handoffs"),
        ("Visibility", "One network view"),
        ("Flexibility", "Respond to peaks faster"),
        ("Scalability", "Grow without full ownership"),
    ]
    left = 0.7
    gap = 0.14
    width = (12.0 - gap * 4) / 5
    for idx, (k, v) in enumerate(pillars):
        x = left + idx * (width + gap)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(strip_y), Inches(width), Inches(1.12))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT_LINE
        add_textbox(slide, Inches(x + 0.12), Inches(strip_y + 0.16), Inches(width - 0.24), Inches(0.22), k,
                    font_size=12.5, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.12), Inches(strip_y + 0.44), Inches(width - 0.24), Inches(0.34), v,
                    font_size=11, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, "The deck uses an OM lens throughout: platform design, standardization, network coordination, and operational trade-offs.", 2)


def slide3(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "Alibaba's scale created the logistics challenge; Cainiao was built to solve it",
              "Alibaba's commerce growth made logistics an operating bottleneck, so Cainiao was founded in 2013 as a data-centered network rather than another asset-heavy courier.")
    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.92), Inches(4.05), Inches(4.95))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = WHITE
    left_panel.line.color.rgb = LIGHT_LINE
    add_section_label(slide, Inches(0.96), Inches(2.12), "TIMELINE + POSITIONING")
    add_textbox(slide, Inches(0.96), Inches(2.42), Inches(3.55), Inches(0.4), "1999 | Alibaba founded", font_size=18, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(0.96), Inches(2.82), Inches(3.45), Inches(0.34), "Large-scale e-commerce created a logistics coordination problem, not just a delivery-volume problem.", font_size=14, color=MUTED)
    add_textbox(slide, Inches(0.96), Inches(3.35), Inches(3.55), Inches(0.4), "2013 | Cainiao founded", font_size=18, color=PRIMARY, bold=True)
    add_textbox(slide, Inches(0.96), Inches(3.75), Inches(3.48), Inches(0.5), "Mission: build a smart logistics network that enables 24-hour delivery in China and 72-hour delivery worldwide.", font_size=14, color=MUTED)
    add_textbox(slide, Inches(0.96), Inches(4.56), Inches(3.42), Inches(0.22), "Why it matters", font_size=12, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.96), Inches(4.82), Inches(3.5), Inches(0.84), "Cainiao behaves more like a logistics operating system: it sets standards, shares data, and orchestrates partners across the network.", font_size=16, color=TEXT)
    cards = [
        ("Vision", "24h in China\n72h worldwide"),
        ("Reach", "200+ countries\nand regions"),
        ("FY2025 revenue", "RMB 101.272B"),
        ("Adj. EBITA", "RMB 302M"),
        ("Infrastructure", "1,100+ warehouses\n170,000 stations"),
    ]
    positions = [
        (5.12, 2.02), (8.72, 2.02),
        (5.12, 3.72), (8.72, 3.72),
        (5.12, 5.42)
    ]
    for (label, value), (x, y) in zip(cards, positions):
        add_card(slide, Inches(x), Inches(y), Inches(3.08), Inches(1.35), label, value,
                 fill=WHITE, text_color=PRIMARY)
    add_footer(slide, "Source: Cainiao official website; Alibaba Group FY2025 annual report.", 3)


def slide4(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "The operational problem is coordination loss across a fragmented network",
              "Parcel demand is large, volatile, and fragmented across many actors. Without shared standards and visibility, local delays become system-wide inefficiency.")
    center = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(3.0), Inches(3.2), Inches(1.05))
    center.fill.solid()
    center.fill.fore_color.rgb = PRIMARY
    center.line.color.rgb = PRIMARY
    add_textbox(slide, Inches(5.18), Inches(3.22), Inches(2.82), Inches(0.5),
                "E-commerce demand shocks expose coordination failures", font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    pain_points = [
        (1.0, 2.1, "Volatile peaks", "Shopping festivals and promotions amplify demand spikes."),
        (9.3, 2.1, "Fragmented carriers", "Different firms operate with different rules and systems."),
        (1.0, 4.65, "Poor end-to-end visibility", "Exceptions are discovered late and handled manually."),
        (9.3, 4.65, "Data inconsistency", "Address, parcel, and status data are not standardized."),
        (5.05, 5.68, "High coordination cost", "More handoffs mean more waiting, rework, and customer uncertainty."),
    ]
    for x, y, title, body in pain_points:
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.65), Inches(1.05))
        box.fill.solid()
        box.fill.fore_color.rgb = BG
        box.line.color.rgb = LIGHT_LINE
        add_textbox(slide, Inches(x + 0.12), Inches(y + 0.12), Inches(2.35), Inches(0.18), title,
                    font_size=11.5, color=PRIMARY, bold=True)
        add_textbox(slide, Inches(x + 0.12), Inches(y + 0.34), Inches(2.35), Inches(0.48), body,
                    font_size=10.2, color=MUTED)
    for line in [
        (Inches(3.65), Inches(2.62), Inches(5.0), Inches(3.22)),
        (Inches(8.2), Inches(2.62), Inches(8.95), Inches(3.22)),
        (Inches(3.65), Inches(5.0), Inches(5.0), Inches(4.02)),
        (Inches(8.2), Inches(5.0), Inches(8.95), Inches(4.02)),
    ]:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, *line)
        conn.line.color.rgb = LIGHT_LINE
        conn.line.width = Pt(1.6)
    add_footer(slide, "Operations implication: the constraint is not only physical capacity; it is network coordination under uncertainty.", 4)


def slide5(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "The traditional model is linear, siloed, and slow to react",
              "When each actor optimizes its own segment, the end-to-end network becomes harder to coordinate and exceptions are harder to resolve.")
    flow_labels = ["Seller", "Courier", "Warehouse", "National hub", "Last mile", "Customer"]
    x_positions = [0.9, 2.55, 4.35, 6.25, 8.35, 10.45]
    for idx, (label, x) in enumerate(zip(flow_labels, x_positions)):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(3.1), Inches(1.35), Inches(0.72))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT_LINE
        add_textbox(slide, Inches(x + 0.08), Inches(3.26), Inches(1.18), Inches(0.18), label,
                    font_size=12, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
        if idx < len(flow_labels) - 1:
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.35), Inches(3.46), Inches(x_positions[idx + 1]), Inches(3.46))
            conn.line.color.rgb = ACCENT
            conn.line.width = Pt(2.0)
    friction = [
        (1.92, "Manual handoff"),
        (3.85, "No shared data"),
        (5.68, "Delayed exceptions"),
        (7.72, "Low ETA accuracy"),
        (9.8, "Local optimization"),
    ]
    for x, txt in friction:
        add_chip(slide, Inches(x), Inches(2.45), Inches(1.3), txt)
    right_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(4.55), Inches(3.05), Inches(1.55))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = WHITE
    right_box.line.color.rgb = LIGHT_LINE
    add_textbox(slide, Inches(9.22), Inches(4.74), Inches(2.62), Inches(0.2), "OM consequence", font_size=11, color=ACCENT, bold=True)
    add_bullets(slide, Inches(9.18), Inches(5.0), Inches(2.72), Inches(0.92),
                ["Buffers grow to absorb uncertainty", "Exception handling stays manual", "Service quality varies across nodes"],
                font_size=11.5)
    add_textbox(slide, Inches(0.92), Inches(5.05), Inches(6.8), Inches(0.7),
                "Linear ownership can optimize each segment, but it does not create a shared network view. That is why scale alone does not solve the coordination problem.",
                font_size=18, color=TEXT)
    add_footer(slide, "The pain point is structural: information moves slower than parcels, so the chain reacts too late.", 5)


def slide6(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Cainiao's platform layer lets information move before goods do",
              "The platform connects sellers, warehouses, carriers, customs, pickup points, and customers through one shared data and coordination layer.")
    hub = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(2.75), Inches(2.9), Inches(1.1))
    hub.fill.solid()
    hub.fill.fore_color.rgb = PRIMARY
    hub.line.color.rgb = PRIMARY
    add_textbox(slide, Inches(5.42), Inches(2.98), Inches(2.46), Inches(0.46),
                "Cainiao data & coordination platform", font_size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    nodes = [
        (1.1, 2.1, "Sellers"),
        (1.3, 4.95, "Warehouses"),
        (4.0, 5.65, "Line-haul carriers"),
        (8.95, 5.0, "Pickup / drop-off"),
        (9.3, 2.2, "Customs / cross-border"),
        (4.2, 1.45, "Consumers"),
    ]
    for x, y, label in nodes:
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(1.62), Inches(0.72))
        s.fill.solid()
        s.fill.fore_color.rgb = BG
        s.line.color.rgb = LIGHT_LINE
        add_textbox(slide, Inches(x + 0.08), Inches(y + 0.2), Inches(1.46), Inches(0.16), label,
                    font_size=11.5, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 0.81), Inches(y + 0.36), Inches(6.65), Inches(3.3))
        conn.line.color.rgb = ACCENT
        conn.line.width = Pt(1.6)
    for idx, text in enumerate(["Standards", "Visibility", "Routing", "Forecasting"]):
        add_chip(slide, Inches(4.72 + idx * 1.45), Inches(4.32), Inches(1.22), text,
                 fill=ACCENT_SOFT if idx % 2 == 0 else WHITE)
    add_textbox(slide, Inches(0.92), Inches(6.32), Inches(10.8), Inches(0.34),
                "Key line: Cainiao moves information so goods can move faster.", font_size=16.5, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Source: Cainiao official website and product descriptions.", 6)


def slide7(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "Standardized data is Cainiao's first operating lever",
              "Before the network can optimize routes or warehouse flows, it needs common parcel data, common addresses, and a shared control view.")
    stack_x = 0.95
    sources = ["Merchant order data", "Carrier status updates", "Warehouse events", "Customer-facing tracking"]
    for idx, label in enumerate(sources):
        add_card(slide, Inches(stack_x), Inches(2.15 + idx * 0.72), Inches(2.48), Inches(0.56), label, "",
                 fill=WHITE, text_color=PRIMARY)
    hub = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.1), Inches(3.0), Inches(3.0), Inches(1.0))
    hub.fill.solid()
    hub.fill.fore_color.rgb = PRIMARY
    hub.line.color.rgb = PRIMARY
    add_textbox(slide, Inches(4.32), Inches(3.25), Inches(2.56), Inches(0.34), "Standardized data hub", font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for idx in range(len(sources)):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.43), Inches(2.43 + idx * 0.72), Inches(4.1), Inches(3.5))
        conn.line.color.rgb = LIGHT_LINE
    outputs = ["Unified e-waybill", "Address normalization", "Control-tower visibility"]
    for idx, label in enumerate(outputs):
        add_card(slide, Inches(8.1), Inches(2.4 + idx * 1.0), Inches(3.35), Inches(0.74), label,
                 "Less ambiguity, fewer coordination losses." if idx == 0 else
                 "One version of parcel status across the network." if idx == 2 else
                 "Cleaner parcel identities for routing and handoff decisions.",
                 fill=WHITE, text_color=PRIMARY)
    add_textbox(slide, Inches(0.95), Inches(5.85), Inches(3.45), Inches(0.62),
                "OM implication: standardization reduces uncertainty before any optimization algorithm starts.", font_size=17, color=TEXT)
    add_bullets(slide, Inches(8.05), Inches(5.45), Inches(3.4), Inches(0.95),
                ["Shared formats cut reconciliation work", "Visibility improves exception response", "The data layer becomes an operating asset"], font_size=11.3)
    add_footer(slide, "Source: Cainiao official site and official product/technology materials.", 7)


def slide8(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "AI routing turns parcel scale into better allocation decisions",
              "The point of AI here is not novelty - it is faster, more accurate operational choices across routing, sorting, and delivery-zone assignment.")
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.88), Inches(2.0), Inches(6.15), Inches(4.45))
    panel.fill.solid()
    panel.fill.fore_color.rgb = SECONDARY
    panel.line.color.rgb = SECONDARY
    add_textbox(slide, Inches(1.15), Inches(2.2), Inches(5.6), Inches(0.24), "Illustrative routing dashboard", font_size=11, color=RGBColor(0xB8, 0xC5, 0xD3), bold=True)
    nodes = [(1.55, 4.2), (2.3, 3.0), (3.7, 3.55), (4.8, 2.6), (5.9, 4.15)]
    for x, y in nodes:
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.24), Inches(0.24))
        circ.fill.solid()
        circ.fill.fore_color.rgb = ACCENT
        circ.line.color.rgb = ACCENT
    connections = [
        ((1.67, 4.32), (2.42, 3.12)),
        ((2.42, 3.12), (3.82, 3.67)),
        ((3.82, 3.67), (4.92, 2.72)),
        ((4.92, 2.72), (6.02, 4.27)),
        ((2.42, 3.12), (5.98, 4.25)),
    ]
    for (x1, y1), (x2, y2) in connections:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = RGBColor(0xF9, 0xAA, 0x54)
        conn.line.width = Pt(2.0)
    add_metric_band(
        slide,
        [
            ("AI zone-code accuracy", "98–99%", "accent"),
            ("Daily volume context", "Tens of millions", "default"),
            ("Labor impact", "3–5% lower labor cost", "accent"),
        ],
        top=5.75,
    )
    add_textbox(slide, Inches(7.45), Inches(2.1), Inches(4.35), Inches(0.86),
                "What improves operationally", font_size=12, color=ACCENT, bold=True)
    add_bullets(slide, Inches(7.38), Inches(2.5), Inches(4.2), Inches(1.32),
                ["Route decisions update with live conditions", "Sorting and dispatch use better parcel grouping", "Delivery zones become more consistent and less labor-intensive"],
                font_size=14)
    add_textbox(slide, Inches(7.45), Inches(4.28), Inches(4.15), Inches(1.2),
                "The OM win is not just speed. Better allocation reduces avoidable labor, improves flow stability, and makes high parcel volume manageable without proportionally increasing coordination effort.",
                font_size=16, color=TEXT)
    add_footer(slide, "Source: INFORMS practice summary on AI-generated delivery zone codes.", 8)


def slide9(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "Smart warehouses and pre-positioned inventory compress lead time",
              "Warehouse automation and network design work together: the facility reduces processing time, while inventory pre-positioning reduces travel time.")
    add_textbox(slide, Inches(0.88), Inches(2.1), Inches(3.48), Inches(1.12),
                "Cainiao reports 1,100+ warehouses totaling about 16.5 million sqm, supported by 380 sorting centers. The operating idea is simple: place inventory closer to likely demand and process parcels faster once they arrive.", font_size=17)
    add_metric_band(
        slide,
        [
            ("Warehouses", "1,100+", "accent"),
            ("Floor area", "≈16.5M sqm", "default"),
            ("Sorting centers", "380", "accent"),
        ],
        top=3.55,
    )
    warehouse = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.35), Inches(2.05), Inches(6.35), Inches(3.95))
    warehouse.fill.solid()
    warehouse.fill.fore_color.rgb = WHITE
    warehouse.line.color.rgb = LIGHT_LINE
    add_textbox(slide, Inches(5.6), Inches(2.22), Inches(5.7), Inches(0.2), "Illustrative warehouse operating view", font_size=11, color=MUTED, bold=True)
    rack_xs = [5.9, 6.95, 8.0, 9.05]
    for x in rack_xs:
        rack = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(2.7), Inches(0.58), Inches(1.62))
        rack.fill.solid()
        rack.fill.fore_color.rgb = RGBColor(0xCF, 0xD9, 0xE3)
        rack.line.color.rgb = RGBColor(0xA3, 0xB4, 0xC5)
        for shelf in [3.05, 3.42, 3.79]:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(shelf), Inches(x + 0.58), Inches(shelf))
            line.line.color.rgb = WHITE
    for x, y in [(6.0, 4.8), (7.25, 4.55), (8.35, 4.95)]:
        agv = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.6), Inches(0.3))
        agv.fill.solid()
        agv.fill.fore_color.rgb = ACCENT
        agv.line.color.rgb = ACCENT
    for (x1, y1), (x2, y2) in [((6.3, 4.95), (7.55, 4.7)), ((7.55, 4.7), (8.65, 5.1)), ((8.65, 5.1), (10.4, 5.1))]:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = ACCENT
        conn.line.width = Pt(1.8)
    add_card(slide, Inches(0.88), Inches(5.0), Inches(2.2), Inches(0.88), "Network logic", "Forecast -> pre-position -> pick/pack -> last mile",
             fill=WHITE, text_color=PRIMARY)
    add_card(slide, Inches(3.18), Inches(5.0), Inches(1.95), Inches(0.88), "OM benefit", "Less travel + faster processing",
             fill=WHITE, text_color=PRIMARY)
    add_footer(slide, "Source: Cainiao official site.", 9)


def slide10(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Platform governance scales the network without owning every asset",
              "The model works only if partners can plug into shared rules, shared interfaces, and a common performance discipline.")
    hub = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(4.95), Inches(2.65), Inches(2.1), Inches(2.1))
    hub.fill.solid()
    hub.fill.fore_color.rgb = PRIMARY
    hub.line.color.rgb = PRIMARY
    add_textbox(slide, Inches(5.18), Inches(3.3), Inches(1.65), Inches(0.44), "Cainiao governance layer", font_size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    ring_items = [
        (2.0, 2.1, "Service standards"),
        (8.8, 2.1, "Shared APIs & access"),
        (2.0, 5.05, "Exception escalation"),
        (8.8, 5.05, "Performance monitoring"),
    ]
    for x, y, label in ring_items:
        item = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.4), Inches(0.8))
        item.fill.solid()
        item.fill.fore_color.rgb = BG
        item.line.color.rgb = LIGHT_LINE
        add_textbox(slide, Inches(x + 0.08), Inches(y + 0.22), Inches(2.24), Inches(0.18), label,
                    font_size=12.2, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 1.2), Inches(y + 0.4), Inches(6.0), Inches(3.7))
        conn.line.color.rgb = ACCENT
        conn.line.width = Pt(1.4)
    add_chip(slide, Inches(9.15), Inches(1.46), Inches(2.15), "170,000 pickup / drop-off stations")
    add_textbox(slide, Inches(0.9), Inches(6.0), Inches(10.9), Inches(0.38),
                "Scale comes from orchestration: partners keep operating the physical network, while Cainiao improves the rules, data, and decisions that coordinate that network.", font_size=17.5, color=TEXT)
    add_footer(slide, "Source: Cainiao official site.", 10)


def slide11(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "Results point to stronger reach, visibility, and economic viability",
              "The evidence supports a more scalable operating model, even if not every outcome can be reduced to a single cost or speed claim.")
    dashboard = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.82), Inches(2.0), Inches(11.0), Inches(4.75))
    dashboard.fill.solid()
    dashboard.fill.fore_color.rgb = WHITE
    dashboard.line.color.rgb = LIGHT_LINE
    metrics = [
        ("Service promise", "24h China / 72h worldwide"),
        ("Global reach", "200+ countries & regions"),
        ("PUDO network", "170,000 stations"),
        ("FY2025 revenue", "RMB 101.272B"),
        ("Adj. EBITA", "RMB 302M"),
        ("Routing proof point", "98–99% AI zone-code accuracy"),
    ]
    positions = [(1.05, 2.32), (4.12, 2.32), (7.19, 2.32), (1.05, 4.05), (4.12, 4.05), (7.19, 4.05)]
    for (label, value), (x, y) in zip(metrics, positions):
        add_card(slide, Inches(x), Inches(y), Inches(2.68), Inches(1.35), label, value,
                 fill=BG if label in {"Service promise", "Adj. EBITA"} else WHITE)
    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.92), Inches(5.63), Inches(1.55), Inches(0.6))
    note.fill.solid()
    note.fill.fore_color.rgb = ACCENT_SOFT
    note.line.color.rgb = ACCENT
    add_textbox(slide, Inches(10.02), Inches(5.78), Inches(1.34), Inches(0.16), "Global 5-day / 10-day services", font_size=9.2, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Source: Cainiao official site; Alibaba FY2025 annual report; INFORMS practice summary.", 11)


def slide12(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, WHITE)
    add_title(slide, "Cainiao wins by synchronizing information before parcels move",
              "The difference is not simply technology adoption; it is a different operating logic for how the network is designed and governed.")
    table_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(2.0), Inches(10.95), Inches(4.95))
    table_bg.fill.solid()
    table_bg.fill.fore_color.rgb = BG
    table_bg.line.color.rgb = LIGHT_LINE
    headers = [("Dimension", 1.0), ("Traditional logistics", 3.65), ("Cainiao model", 7.45)]
    for label, x in headers:
        add_textbox(slide, Inches(x), Inches(2.24), Inches(2.9), Inches(0.22), label, font_size=12.5, color=PRIMARY, bold=True)
    rows = [
        ("Asset model", "Owns more of the network directly", "Uses an asset-light platform to coordinate others"),
        ("Information flow", "Local and delayed", "Shared and closer to real time"),
        ("Coordination logic", "Sequential handoffs", "Network orchestration through standards and data"),
        ("Scalability", "Capacity grows with owned assets", "Capacity grows through partner participation"),
        ("Role of data", "Supports reporting", "Drives routing, visibility, and control"),
    ]
    y = 2.72
    for idx, (dim, trad, cai) in enumerate(rows):
        row_fill = WHITE if idx % 2 == 0 else RGBColor(0xEE, 0xF3, 0xF7)
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(y), Inches(10.45), Inches(0.78))
        band.fill.solid()
        band.fill.fore_color.rgb = row_fill
        band.line.color.rgb = BG
        add_textbox(slide, Inches(1.12), Inches(y + 0.18), Inches(2.1), Inches(0.22), dim, font_size=11.5, color=PRIMARY, bold=True)
        add_textbox(slide, Inches(3.72), Inches(y + 0.12), Inches(2.95), Inches(0.42), trad, font_size=11.2, color=TEXT)
        add_textbox(slide, Inches(7.52), Inches(y + 0.12), Inches(3.55), Inches(0.42), cai, font_size=11.2, color=TEXT)
        y += 0.78
    add_footer(slide, "This comparison is analytical rather than brand-centric: the operating logic changes before the network outcomes do.", 12)


def slide13(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG)
    add_title(slide, "The model is powerful, but it trades ownership for dependency",
              "A platform model reduces capital intensity, but it also increases exposure to partner quality, governance complexity, and data/compliance risk.")
    risks = [
        ("Partner dependence", "Performance still depends on the reliability and incentives of outside carriers and local operators."),
        ("Data quality & compliance", "If parcel identities, addresses, or privacy controls are weak, the platform becomes harder to trust and govern."),
        ("International complexity", "Cross-border operations multiply customs, regulation, and service-consistency challenges."),
        ("Margin pressure", "Coordinating the ecosystem can improve scale, but it does not remove cost pressure or strategic dependence on Alibaba."),
    ]
    coords = [(0.95, 2.15), (6.65, 2.15), (0.95, 4.42), (6.65, 4.42)]
    fills = [WHITE, RGBColor(0xFF, 0xF6, 0xE9), WHITE, RGBColor(0xF2, 0xF7, 0xFB)]
    for (title, body), (x, y), fill in zip(risks, coords, fills):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.02), Inches(1.62))
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = LIGHT_LINE
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.15), Inches(4.7), Inches(0.22), title,
                    font_size=12.5, color=PRIMARY, bold=True)
        add_textbox(slide, Inches(x + 0.15), Inches(y + 0.45), Inches(4.68), Inches(0.86), body,
                    font_size=13.2, color=TEXT)
    add_textbox(slide, Inches(0.98), Inches(6.45), Inches(10.8), Inches(0.3),
                "Management implication: algorithms matter, but governance is what keeps the platform dependable at scale.", font_size=16.5, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Critical evaluation is part of the case: the platform is scalable because governance keeps the ecosystem coordinated.", 13)


def slide14(prs: Presentation, hero: Path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if hero.exists():
        slide.shapes.add_picture(str(hero), 0, 0, width=prs.slide_width, height=prs.slide_height)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = PRIMARY
    overlay.fill.transparency = 0.18
    overlay.line.color.rgb = PRIMARY
    add_title(slide, "Four operations-management lessons from Cainiao", None, dark=True)
    lessons = [
        ("Standardize data first", "Shared formats reduce uncertainty before execution begins."),
        ("Optimize decisions dynamically", "Routing and allocation improve when the system sees the network as a whole."),
        ("Design network and inventory together", "Lead time falls when placement and processing are coordinated."),
        ("Govern partners like an operating system", "Scale comes from orchestration, not full ownership."),
    ]
    positions = [(0.9, 2.18), (6.7, 2.18), (0.9, 4.12), (6.7, 4.12)]
    for (title, body), (x, y) in zip(lessons, positions):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.1), Inches(1.3))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.fill.transparency = 0.08
        box.line.color.rgb = RGBColor(0xDB, 0xE5, 0xEE)
        add_textbox(slide, Inches(x + 0.14), Inches(y + 0.14), Inches(4.8), Inches(0.22), title,
                    font_size=12.5, color=PRIMARY, bold=True)
        add_textbox(slide, Inches(x + 0.14), Inches(y + 0.46), Inches(4.78), Inches(0.5), body,
                    font_size=13, color=TEXT)
    add_textbox(slide, Inches(0.82), Inches(6.4), Inches(11.1), Inches(0.32),
                "Conclusion: Cainiao transforms logistics from a sequence of handoffs into a data-driven operating system.", font_size=18.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "Vertex-assisted content supplement and hero visual; final deck rebuilt for an Operations Management audience.", 14, dark=True)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    hero = copy_cover_image()
    _ = load_vertex_spec()
    slide1(prs, hero)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    slide6(prs)
    slide7(prs)
    slide8(prs)
    slide9(prs)
    slide10(prs)
    slide11(prs)
    slide12(prs)
    slide13(prs)
    slide14(prs, hero)
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    build_deck()
